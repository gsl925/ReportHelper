import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, PanedWindow
import pytesseract
from PIL import Image, ImageGrab
import os
import re
import pyperclip
from pptx import Presentation
from pptx.util import Inches, Pt
import extract_msg
import sys # 引入 sys 模組以處理打包後的路徑

# 引入拖曳功能的核心函式庫
from tkinterdnd2 import DND_FILES, TkinterDnD

# --- 設定 Tesseract 路徑 (僅 Windows 需要) ---
# if os.name == 'nt':
#     pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# --- 定義固定的 PowerPoint 檔案名稱 ---
MASTER_PPTX_FILENAME = "Weekly Report_JimChuang.pptx"

# --- PROMPTS 已被移除，將改為從外部檔案讀取 ---

class ReportHelperApp_v15_3:
    def __init__(self, root):
        self.root = root
        
        # --- 新增：讀取外部 Prompt 檔案的邏輯 ---
        try:
            # 決定檔案路徑的基準點，使其在 .py 和 .exe 環境下都能運作
            if getattr(sys, 'frozen', False):
                # 如果是打包後的 .exe
                base_path = os.path.dirname(sys.executable)
            else:
                # 如果是直接執行 .py 檔
                base_path = os.path.dirname(__file__)

            with open(os.path.join(base_path, 'prompt_single_issue.txt'), 'r', encoding='utf-8') as f:
                self.single_issue_prompt = f.read()
            with open(os.path.join(base_path, 'prompt_multi_issue.txt'), 'r', encoding='utf-8') as f:
                self.multi_issue_prompt = f.read()
        except FileNotFoundError as e:
            messagebox.showerror("錯誤：找不到 Prompt 檔案", f"找不到必要的設定檔：\n{os.path.basename(e.filename)}\n\n請確認 prompt_single_issue.txt 和 prompt_multi_issue.txt 檔案與主程式放在同一個資料夾中。")
            self.root.destroy() # 找不到檔案就直接關閉程式
            return # 結束初始化
        except Exception as e:
            messagebox.showerror("錯誤：讀取 Prompt 檔案失敗", f"讀取設定檔時發生錯誤：\n{e}")
            self.root.destroy()
            return

        self.root.title(f"報告整理小幫手 v15.3 (進階分析版)")
        self.root.geometry("1200x700")

        paned_window = PanedWindow(root, orient=tk.HORIZONTAL, sashrelief=tk.RAISED)
        paned_window.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

        self.left_frame = tk.Frame(paned_window, padx=5)
        self.create_left_panel(self.left_frame)
        paned_window.add(self.left_frame, stretch="always")

        right_frame = tk.Frame(paned_window, padx=5)
        self.create_right_panel(right_frame)
        paned_window.add(right_frame, stretch="always")

        self.root.drop_target_register(DND_FILES)
        self.root.dnd_bind('<<Drop>>', self.handle_drop)
        
        self.original_bg = self.left_frame.cget("background")
        self.root.dnd_bind('<<DragEnter>>', self.on_drag_enter)
        self.root.dnd_bind('<<DragLeave>>', self.on_drag_leave)

    def on_drag_enter(self, event):
        self.left_frame.config(bg="#E0E0E0")
        return event.action

    def on_drag_leave(self, event):
        self.left_frame.config(bg=self.original_bg)

    def handle_drop(self, event):
        self.on_drag_leave(event)
        filepaths_str = event.data
        filepaths = self.root.tk.splitlist(filepaths_str)
        if not filepaths:
            return
        
        self.process_file_list(filepaths)

    def create_left_panel(self, parent):
        controls_frame = tk.Frame(parent)
        controls_frame.pack(fill=tk.X, pady=(0, 10))
        
        project_frame = tk.Frame(controls_frame)
        project_frame.pack(fill=tk.X, pady=(0, 10))
        project_label = tk.Label(project_frame, text="專案名稱 (選填):")
        project_label.pack(side=tk.LEFT, padx=(0, 5))
        self.project_name_entry = tk.Entry(project_frame)
        self.project_name_entry.pack(side=tk.LEFT, fill=tk.X, expand=True)

        upload_paste_frame = tk.Frame(controls_frame)
        upload_paste_frame.pack(fill=tk.X)
        upload_button = tk.Button(upload_paste_frame, text="1. 上傳檔案 (可多選)", command=self.upload_file)
        upload_button.pack(side=tk.LEFT, padx=(0, 5))
        
        paste_button = tk.Button(upload_paste_frame, text="2. 從剪貼簿附加圖片", command=self.paste_from_clipboard)
        paste_button.pack(side=tk.LEFT, padx=5)
        
        self.status_label = tk.Label(upload_paste_frame, text="請上傳或拖曳檔案至此...", fg="blue")
        self.status_label.pack(side=tk.LEFT, padx=10)

        text_frame = tk.LabelFrame(parent, text="步驟 A: 辨識結果 (原始文字)", padx=5, pady=5)
        text_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        self.text_area = scrolledtext.ScrolledText(text_frame, wrap=tk.WORD)
        self.text_area.pack(fill=tk.BOTH, expand=True)
        
        button_frame = tk.Frame(parent)
        button_frame.pack(fill=tk.X)
        
        single_button = tk.Button(button_frame, text="步驟 B: 分析為「單一問題」", command=self.copy_single_issue_prompt, bg="#2196F3", fg="white", height=2)
        single_button.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 5))
        
        multi_button = tk.Button(button_frame, text="步驟 B: 分析為「多個問題」", command=self.copy_multi_issue_prompt, bg="#4CAF50", fg="white", height=2)
        multi_button.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(5, 0))

    def create_right_panel(self, parent):
        genai_frame = tk.LabelFrame(parent, text="步驟 C: 在此貼上 GenAI 產生的報告", padx=5, pady=5)
        genai_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        self.genai_output_area = scrolledtext.ScrolledText(genai_frame, wrap=tk.WORD)
        self.genai_output_area.pack(fill=tk.BOTH, expand=True)

        ppt_button_text = f"步驟 D: 將報告新增至彙總簡報 ({MASTER_PPTX_FILENAME})"
        ppt_button = tk.Button(parent, text=ppt_button_text, command=self.add_to_master_ppt, bg="#E91E63", fg="white", height=2, font=("Arial", 10, "bold"))
        ppt_button.pack(fill=tk.X)

    def _prepare_prompt(self, base_prompt):
        content = self.text_area.get("1.0", tk.END).strip()
        if not content:
            messagebox.showwarning("內容為空", "辨識結果為空，無法複製。")
            return None
        
        project_name = self.project_name_entry.get().strip()
        project_info_text = f"The user has specified the project name is: '{project_name}'." if project_name else "The user did not specify a project name."
        
        final_prompt = base_prompt.replace("{PROJECT_NAME_HOLDER}", project_info_text)
        return f"{final_prompt}\n\n{content}"

    def copy_single_issue_prompt(self):
        # 使用 self.single_issue_prompt
        full_prompt = self._prepare_prompt(self.single_issue_prompt)
        if full_prompt:
            pyperclip.copy(full_prompt)
            messagebox.showinfo("複製成功", "「單一問題」分析 Prompt 及內容已複製到剪貼簿！")

    def copy_multi_issue_prompt(self):
        # 使用 self.multi_issue_prompt
        full_prompt = self._prepare_prompt(self.multi_issue_prompt)
        if full_prompt:
            pyperclip.copy(full_prompt)
            messagebox.showinfo("複製成功", "「多個問題」分析 Prompt 及內容已複製到剪貼簿！")

    def upload_file(self):
        filetypes = (
            ("支援的檔案", "*.png *.jpg *.jpeg *.txt *.msg"),
            ("Email 檔案", "*.msg"),
            ("圖片檔案", "*.png *.jpg *.jpeg"), 
            ("文字檔案", "*.txt"), 
            ("所有檔案", "*.*")
        )
        filepaths = filedialog.askopenfilenames(filetypes=filetypes)
        if not filepaths: return
        
        self.process_file_list(filepaths)

    def process_file_list(self, filepaths):
        self.text_area.delete('1.0', tk.END)
        total_files = len(filepaths)

        for i, file_path in enumerate(filepaths):
            if i > 0:
                separator = f"\n\n{'='*20} 檔案 {i+1} {'='*20}\n\n"
                self.text_area.insert(tk.END, separator)

            self.status_label.config(text=f"正在處理第 {i+1}/{total_files} 個檔案: {os.path.basename(file_path)}...", fg="blue")
            self.root.update_idletasks()

            if not os.path.exists(file_path):
                self.text_area.insert(tk.END, f"錯誤：找不到檔案 {file_path}")
                continue

            file_ext = os.path.splitext(file_path)[1].lower()
            try:
                if file_ext in ['.png', '.jpg', '.jpeg']: 
                    self.process_image_object(Image.open(file_path))
                elif file_ext in ['.txt']: 
                    self.process_text_file(file_path)
                elif file_ext in ['.msg']:
                    self.process_msg_file(file_path)
                else: 
                    messagebox.showwarning("不支援的格式", f"不支援的檔案格式: {file_ext}")
            except Exception as e: 
                self.handle_error(e, f"處理檔案 {os.path.basename(file_path)} 時發生錯誤")

        self.status_label.config(text=f"全部 {total_files} 個檔案處理完成！", fg="green")

    def paste_from_clipboard(self):
        try:
            image = ImageGrab.grabclipboard()
            if not isinstance(image, Image.Image):
                messagebox.showinfo("提示", "剪貼簿中沒有圖片。")
                self.status_label.config(text="剪貼簿中沒有圖片。", fg="orange")
                return

            if self.text_area.get("1.0", tk.END).strip():
                separator = f"\n\n{'='*20} 來自剪貼簿的新增圖片 {'='*20}\n\n"
                self.text_area.insert(tk.END, separator)
            
            self.status_label.config(text="正在辨識剪貼簿中的圖片...", fg="blue")
            self.root.update_idletasks()

            self.process_image_object(image)
            
            self.status_label.config(text="已從剪貼簿附加圖片內容！", fg="green")

        except Exception as e: 
            self.handle_error(e, "從剪貼簿讀取圖片時發生錯誤")

    def process_image_object(self, image_obj):
        try:
            lang_models = 'chi_tra+chi_sim+eng'
            extracted_text = pytesseract.image_to_string(image_obj, lang=lang_models)
            self.text_area.insert(tk.END, extracted_text)
        except pytesseract.TesseractNotFoundError:
            messagebox.showerror("Tesseract 未找到", "找不到 Tesseract OCR 引擎。請確認已安裝且路徑正確。")
            self.status_label.config(text="Tesseract 未安裝或路徑錯誤！", fg="red")
        except Exception as e: raise e

    def process_text_file(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as f: content = f.read()
        except Exception:
            try:
                with open(file_path, 'r', encoding='gbk') as f: content = f.read()
            except Exception as e2:
                self.handle_error(e2, "讀取文字檔失敗")
                return
        self.text_area.insert(tk.END, content)

    def process_msg_file(self, file_path):
        try:
            msg = extract_msg.Message(file_path)
            sender = msg.sender
            subject = msg.subject
            body = msg.body

            formatted_content = (
                f"寄件人：{sender}\n"
                f"主旨：{subject}\n\n"
                f"--- 內文 ---\n"
                f"{body}"
            )
            self.text_area.insert(tk.END, formatted_content)
        except Exception as e:
            self.handle_error(e, "解析 Email 檔案時發生錯誤")

    def handle_error(self, e, custom_msg="處理時發生錯誤"):
        messagebox.showerror("錯誤", f"{custom_msg}：\n{e}")
        self.status_label.config(text="處理失敗！", fg="red")

    def add_to_master_ppt(self):
        genai_text = self.genai_output_area.get("1.0", tk.END).strip()
        if not genai_text:
            messagebox.showwarning("內容為空", "請先在右側方塊中貼上 GenAI 產生的報告內容。")
            return

        try:
            # 使用與讀取 Prompt 相同的路徑邏輯，確保能找到 PPT
            if getattr(sys, 'frozen', False):
                base_path = os.path.dirname(sys.executable)
            else:
                base_path = os.path.dirname(__file__)
            pptx_path = os.path.join(base_path, MASTER_PPTX_FILENAME)
        except NameError:
            # 備用方案，以防 __file__ 不可用
            base_path = os.getcwd() 
            pptx_path = os.path.join(base_path, MASTER_PPTX_FILENAME)
            
        try:
            prs = Presentation(pptx_path) if os.path.exists(pptx_path) else Presentation()
            reports = re.split(r'(?=--- 報告 \d+：)', genai_text)
            generated_count = 0
            
            project_name = self.project_name_entry.get().strip()

            for report_text in reports:
                report_text = report_text.strip()
                if not report_text: continue

                match = re.match(r'--- 報告 \d+：(.*?) ---\n(.*)', report_text, re.DOTALL)
                if not match: continue

                ai_title = match.group(1).strip()
                content = match.group(2).strip()

                final_title = f"{project_name} - {ai_title}" if project_name else ai_title

                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = final_title
                
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear() 
                text_frame.word_wrap = True

                for line in content.split('\n'):
                    if not line.strip(): continue 

                    is_indented = line.startswith('  ')
                    level = 1 if is_indented else 0
                    
                    cleaned_line = re.sub(r'^\s*[-*]\s+', '', line).strip()

                    p = text_frame.add_paragraph()
                    p.text = cleaned_line
                    p.level = level

                    if level == 0:
                        p.font.bold = True
                        p.font.size = Pt(18)
                    else:
                        p.font.bold = False
                        p.font.size = Pt(16)
                
                generated_count += 1
            
            if generated_count > 0:
                prs.save(pptx_path)
                messagebox.showinfo("新增成功", f"成功將 {generated_count} 張新的投影片新增至\n'{MASTER_PPTX_FILENAME}'！\n(已套用精簡排版)")
            else:
                messagebox.showwarning("未找到報告", "未能在貼上的內容中找到符合格式的報告，請檢查內容。")

        except PermissionError:
            messagebox.showerror("權限錯誤", f"無法儲存檔案 '{MASTER_PPTX_FILENAME}'。\n\n請先將該 PowerPoint 檔案關閉後再試一次！")
        except Exception as e:
            messagebox.showerror("生成失敗", f"處理 PowerPoint 檔案時發生錯誤：\n{e}")

if __name__ == "__main__":
    root = TkinterDnD.Tk()
    app = ReportHelperApp_v15_3(root)
    root.mainloop()
