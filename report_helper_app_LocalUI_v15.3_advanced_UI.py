# 引入 ttkbootstrap 作為主要的 UI 庫
import ttkbootstrap as ttk
from ttkbootstrap.widgets.scrolled import ScrolledText
from ttkbootstrap.constants import *

# 其他函式庫保持不變
import tkinter as tk
from tkinter import filedialog, messagebox
import pytesseract
from PIL import Image, ImageGrab
import os
import re
import pyperclip
from pptx import Presentation
from pptx.util import Inches, Pt
import extract_msg
import sys

# 引入拖曳功能的核心函式庫
from tkinterdnd2 import DND_FILES, TkinterDnD

# --- 正確的整合方式 ---
class ThemedTkinterDnD(TkinterDnD.Tk):
    def __init__(self, *args, **kwargs):
        themename = kwargs.pop('themename', 'litera')
        super().__init__(*args, **kwargs)
        ttk.Style(theme=themename)

# --- 其餘部分與原程式碼相似 ---
MASTER_PPTX_FILENAME = "Weekly Report_JimChuang.pptx"

class ReportHelperApp_v15_3:
    def __init__(self, root):
        self.root = root
        self.root.title(f"報告整理小幫手 v15.3 (進階分析版)")
        self.root.geometry("1200x750")

        paned_window = ttk.Panedwindow(root, orient=HORIZONTAL)
        paned_window.pack(fill=BOTH, expand=True, padx=10, pady=10)

        self.left_frame = tk.Frame(paned_window) 
        style = ttk.Style.get_instance()
        bg_color = style.colors.get('bg')
        self.left_frame.config(background=bg_color, padx=10, pady=10)

        self.create_left_panel(self.left_frame)
        paned_window.add(self.left_frame, weight=1)

        right_frame = ttk.Frame(paned_window, padding=10)
        self.create_right_panel(right_frame)
        paned_window.add(right_frame, weight=1)

        self.root.drop_target_register(DND_FILES)
        self.root.dnd_bind('<<Drop>>', self.handle_drop)
        
        self.original_bg = str(self.left_frame.cget("background"))
        self.root.dnd_bind('<<DragEnter>>', self.on_drag_enter)
        self.root.dnd_bind('<<DragLeave>>', self.on_drag_leave)
        
        try:
            if getattr(sys, 'frozen', False): base_path = os.path.dirname(sys.executable)
            else: base_path = os.path.dirname(__file__)
            with open(os.path.join(base_path, 'prompt_single_issue.txt'), 'r', encoding='utf-8') as f: self.single_issue_prompt = f.read()
            with open(os.path.join(base_path, 'prompt_multi_issue.txt'), 'r', encoding='utf-8') as f: self.multi_issue_prompt = f.read()
        except FileNotFoundError as e:
            messagebox.showerror("錯誤：找不到 Prompt 檔案", f"找不到必要的設定檔：\n{os.path.basename(e.filename)}\n\n請確認 prompt_single_issue.txt 和 prompt_multi_issue.txt 檔案與主程式放在同一個資料夾中。")
            self.root.destroy()
            return
        except Exception as e:
            messagebox.showerror("錯誤：讀取 Prompt 檔案失敗", f"讀取設定檔時發生錯誤：\n{e}")
            self.root.destroy()
            return

    def create_left_panel(self, parent):
        project_frame = ttk.Frame(parent)
        project_frame.pack(fill=X, pady=(0, 15))
        project_label = ttk.Label(project_frame, text="專案名稱 (選填):")
        project_label.pack(side=LEFT, padx=(0, 10))
        self.project_name_entry = ttk.Entry(project_frame)
        self.project_name_entry.pack(side=LEFT, fill=X, expand=True)

        upload_paste_frame = ttk.Frame(parent)
        upload_paste_frame.pack(fill=X, pady=(0, 10))
        upload_button = ttk.Button(upload_paste_frame, text="1. 上傳檔案", command=self.upload_file, bootstyle="info")
        upload_button.pack(side=LEFT, padx=(0, 5))
        
        paste_button = ttk.Button(upload_paste_frame, text="2. 貼上圖片", command=self.paste_from_clipboard, bootstyle="info-outline")
        paste_button.pack(side=LEFT, padx=5)
        
        self.status_label = ttk.Label(upload_paste_frame, text="請上傳或拖曳檔案至此...", bootstyle="primary")
        self.status_label.pack(side=LEFT, padx=10)

        text_frame = ttk.Labelframe(parent, text="步驟 A: 辨識結果 (原始文字)", padding=5)
        text_frame.pack(fill=BOTH, expand=True, pady=(0, 10))
        self.text_area = ScrolledText(text_frame, wrap=WORD, autohide=True)
        self.text_area.pack(fill=BOTH, expand=True)
        
        button_frame = ttk.Frame(parent)
        button_frame.pack(fill=X, ipady=5)
        
        single_button = ttk.Button(button_frame, text="步驟 B: 分析為「單一問題」", command=self.copy_single_issue_prompt, bootstyle="primary")
        single_button.pack(side=LEFT, fill=X, expand=True, padx=(0, 5))
        
        multi_button = ttk.Button(button_frame, text="步驟 B: 分析為「多個問題」", command=self.copy_multi_issue_prompt, bootstyle="success")
        multi_button.pack(side=LEFT, fill=X, expand=True, padx=(5, 0))

    def create_right_panel(self, parent):
        genai_frame = ttk.Labelframe(parent, text="步驟 C: 在此貼上 GenAI 產生的報告", padding=5)
        genai_frame.pack(fill=BOTH, expand=True, pady=(0, 10))
        self.genai_output_area = ScrolledText(genai_frame, wrap=WORD, autohide=True)
        self.genai_output_area.pack(fill=BOTH, expand=True)

        ppt_button_text = f"步驟 D: 新增至彙總簡報 ({MASTER_PPTX_FILENAME})"
        ppt_button = ttk.Button(parent, text=ppt_button_text, command=self.add_to_master_ppt, bootstyle="danger",)
        ppt_button.pack(fill=X, ipady=5)

    def on_drag_enter(self, event):
        self.left_frame.config(background="#e0e0e0")
        return event.action

    def on_drag_leave(self, event):
        self.left_frame.config(background=self.original_bg)

    def handle_drop(self, event):
        self.on_drag_leave(event)
        filepaths_str = event.data
        filepaths = self.root.tk.splitlist(filepaths_str)
        if not filepaths: return
        self.process_file_list(filepaths)

    def _prepare_prompt(self, base_prompt):
        content = self.text_area.get("1.0", tk.END).strip()
        if not content:
            messagebox.showwarning("內容為空", "辨識結果為空，無法複製。")
            return None
        project_name = self.project_name_entry.get().strip()
        project_info_text = f"The user has specified the project name is: '{project_name}'." if project_name else "The user did not specify a project name."
        final_prompt = base_prompt.replace("{PROJECT_NAME_HOLDER}", project_info_text)
        return f"{final_prompt}\n\n{content}"

    def copy_single_issue_prompt(self):
        full_prompt = self._prepare_prompt(self.single_issue_prompt)
        if full_prompt:
            pyperclip.copy(full_prompt)
            messagebox.showinfo("複製成功", "「單一問題」分析 Prompt 及內容已複製到剪貼簿！")

    def copy_multi_issue_prompt(self):
        full_prompt = self._prepare_prompt(self.multi_issue_prompt)
        if full_prompt:
            pyperclip.copy(full_prompt)
            messagebox.showinfo("複製成功", "「多個問題」分析 Prompt 及內容已複製到剪貼簿！")

    def upload_file(self):
        filetypes = (("支援的檔案", "*.png *.jpg *.jpeg *.txt *.msg"),("Email 檔案", "*.msg"),("圖片檔案", "*.png *.jpg *.jpeg"), ("文字檔案", "*.txt"), ("所有檔案", "*.*"))
        filepaths = filedialog.askopenfilenames(filetypes=filetypes)
        if not filepaths: return
        self.process_file_list(filepaths)

    def process_file_list(self, filepaths):
        self.text_area.delete('1.0', tk.END)
        total_files = len(filepaths)
        for i, file_path in enumerate(filepaths):
            if i > 0:
                separator = f"\n\n{'='*20} 檔案 {i+1} {'='*20}\n\n"
                self.text_area.insert(tk.END, separator)
            self.status_label.config(text=f"正在處理第 {i+1}/{total_files} 個檔案...", bootstyle="info")
            self.root.update_idletasks()
            if not os.path.exists(file_path):
                self.text_area.insert(tk.END, f"錯誤：找不到檔案 {file_path}")
                continue
            file_ext = os.path.splitext(file_path)[1].lower()
            try:
                if file_ext in ['.png', '.jpg', '.jpeg']: self.process_image_object(Image.open(file_path))
                elif file_ext in ['.txt']: self.process_text_file(file_path)
                elif file_ext in ['.msg']: self.process_msg_file(file_path)
                else: messagebox.showwarning("不支援的格式", f"不支援的檔案格式: {file_ext}")
            except Exception as e: self.handle_error(e, f"處理檔案 {os.path.basename(file_path)} 時發生錯誤")
        self.status_label.config(text=f"全部 {total_files} 個檔案處理完成！", bootstyle="success")

    def paste_from_clipboard(self):
        try:
            image = ImageGrab.grabclipboard()
            if not isinstance(image, Image.Image):
                messagebox.showinfo("提示", "剪貼簿中沒有圖片。")
                self.status_label.config(text="剪貼簿中沒有圖片。", bootstyle="warning")
                return
            if self.text_area.get("1.0", tk.END).strip():
                separator = f"\n\n{'='*20} 來自剪貼簿的新增圖片 {'='*20}\n\n"
                self.text_area.insert(tk.END, separator)
            self.status_label.config(text="正在辨識剪貼簿中的圖片...", bootstyle="info")
            self.root.update_idletasks()
            self.process_image_object(image)
            self.status_label.config(text="已從剪貼簿附加圖片內容！", bootstyle="success")
        except Exception as e: self.handle_error(e, "從剪貼簿讀取圖片時發生錯誤")

    def process_image_object(self, image_obj):
        try:
            lang_models = 'chi_tra+chi_sim+eng'
            extracted_text = pytesseract.image_to_string(image_obj, lang=lang_models)
            self.text_area.insert(tk.END, extracted_text)
        except pytesseract.TesseractNotFoundError:
            messagebox.showerror("Tesseract 未找到", "找不到 Tesseract OCR 引擎。請確認已安裝且路徑正確。")
            self.status_label.config(text="Tesseract 未安裝或路徑錯誤！", bootstyle="danger")
        except Exception as e: raise e

    def process_text_file(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as f: content = f.read()
        except Exception:
            try:
                with open(file_path, 'r', encoding='gbk') as f: content = f.read()
            except Exception as e2:
                self.handle_error(e2, "讀取文字檔失敗")
                return
        self.text_area.insert(tk.END, content)

    def process_msg_file(self, file_path):
        try:
            msg = extract_msg.Message(file_path)
            sender, subject, body = msg.sender, msg.subject, msg.body
            formatted_content = (f"寄件人：{sender}\n主旨：{subject}\n\n--- 內文 ---\n{body}")
            self.text_area.insert(tk.END, formatted_content)
        except Exception as e: self.handle_error(e, "解析 Email 檔案時發生錯誤")

    def handle_error(self, e, custom_msg="處理時發生錯誤"):
        messagebox.showerror("錯誤", f"{custom_msg}：\n{e}")
        self.status_label.config(text="處理失敗！", bootstyle="danger")
        
    def add_to_master_ppt(self):
        genai_text = self.genai_output_area.get("1.0", tk.END).strip()
        if not genai_text:
            messagebox.showwarning("內容為空", "請先在右側方塊中貼上 GenAI 產生的報告內容。")
            return
        try:
            if getattr(sys, 'frozen', False): base_path = os.path.dirname(sys.executable)
            else: base_path = os.path.dirname(__file__)
            pptx_path = os.path.join(base_path, MASTER_PPTX_FILENAME)
        except NameError:
            base_path = os.getcwd() 
            pptx_path = os.path.join(base_path, MASTER_PPTX_FILENAME)
        try:
            prs = Presentation(pptx_path) if os.path.exists(pptx_path) else Presentation()
            reports = re.split(r'(?=--- 報告 \d+：)', genai_text)
            generated_count = 0
            project_name = self.project_name_entry.get().strip()
            
            STAR_KEYWORDS = ("情境", "任務", "行動", "結果")
            
            for report_text in reports:
                report_text = report_text.strip()
                if not report_text: continue
                match = re.match(r'--- 報告 \d+：(.*?) ---\n(.*)', report_text, re.DOTALL)
                if not match: continue
                ai_title = match.group(1).strip()
                content = match.group(2).strip()
                final_title = f"{project_name} - {ai_title}" if project_name else ai_title
                slide_layout = prs.slide_layouts[1] # 使用「標題及內容」版面配置
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = final_title
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear() 
                text_frame.word_wrap = True

                # --- 全新修正的階層處理邏輯 ---
                for line in content.split('\n'):
                    original_line = line.strip()
                    if not original_line:
                        continue

                    # 1. 強化清理：移除所有前導符號 (-, *, •) 和頭尾多餘的星號/空白
                    #    例如 "*情境 (Situation)**" -> "情境 (Situation)"
                    cleaned_text = re.sub(r'^\s*[-*•\s`]+', '', original_line) # 移除前導符號
                    cleaned_text = cleaned_text.strip().strip('*').strip() # 移除包圍的星號和空白

                    # 如果清理後為空（例如該行為 "---"），則跳過
                    if not cleaned_text:
                        continue

                    # 2. 判斷清理後的文字是否為 STAR 關鍵字標題
                    is_star_heading = cleaned_text.startswith(STAR_KEYWORDS)

                    if is_star_heading:
                        # 這是主標題 (Level 0)
                        p = text_frame.add_paragraph()
                        # 顯示清理後的文字，例如 "情境 (Situation)"
                        p.text = cleaned_text 
                        p.level = 0
                        p.font.bold = True
                        p.font.size = Pt(18)
                    else:
                        # 這是一般內容，做為子項目 (Level 1)
                        p = text_frame.add_paragraph()
                        # 同樣顯示清理後的文字，以移除項目符號
                        p.text = cleaned_text 
                        p.level = 1
                        p.font.bold = False
                        p.font.size = Pt(16)
                # --- 邏輯修正結束 ---
                generated_count += 1
            
            if generated_count > 0:
                prs.save(pptx_path)
                messagebox.showinfo("新增成功", f"成功將 {generated_count} 張新的投影片新增至\n'{MASTER_PPTX_FILENAME}'！\n(已套用全新階層排版)")
            else:
                messagebox.showwarning("未找到報告", "未能在貼上的內容中找到符合格式的報告，請檢查內容。")
        except PermissionError:
            messagebox.showerror("權限錯誤", f"無法儲存檔案 '{MASTER_PPTX_FILENAME}'。\n\n請先將該 PowerPoint 檔案關閉後再試一次！")
        except Exception as e:
            messagebox.showerror("生成失敗", f"處理 PowerPoint 檔案時發生錯誤：\n{e}")

if __name__ == "__main__":
    # 使用我們自訂的整合類別，並選擇一個主題
    # 其他主題選項: "cosmo", "flatly", "journal", "superhero", "darkly", "cyborg"
    root = ThemedTkinterDnD(themename="litera")
    app = ReportHelperApp_v15_3(root)
    root.mainloop()
