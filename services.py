# services.py
# 核心服務層，處理業務邏輯

import os
import re
import requests
import json
import pytesseract
from PIL import Image
import extract_msg
from pptx import Presentation
from pptx.util import Inches, Pt

from config import OCR_LANGUAGES, STAR_KEYWORDS

class FileProcessorService:
    """處理檔案讀取與文字辨識"""
    def process_image_object(self, image_obj):
        try:
            return pytesseract.image_to_string(image_obj, lang=OCR_LANGUAGES)
        except pytesseract.TesseractNotFoundError:
            raise Exception("找不到 Tesseract OCR 引擎。\n請確認已安裝且路徑設定正確。")
        except Exception as e:
            raise Exception(f"圖片 OCR 辨識失敗：{e}")

    def process_text_file(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()

    def process_msg_file(self, file_path):
        msg = extract_msg.Message(file_path)
        return f"寄件人：{msg.sender}\n主旨：{msg.subject}\n\n--- 內文 ---\n{msg.body}"

class OllamaService:
    """與 Ollama API 進行通訊"""
    def __init__(self, api_url, model):
        self.api_url = api_url
        self.model = model

    def generate(self, prompt):
        try:
            payload = {
                "model": self.model,
                "prompt": prompt,
                "stream": False,
                "options": {
                    "temperature": 0.1,  # 降低溫度，讓輸出更穩定
                    "num_ctx": 4096      # 增加上下文視窗
                }
            }
            response = requests.post(self.api_url, json=payload, timeout=300)
            response.raise_for_status()
            
            response_data = response.json()
            return response_data.get('response', '').strip()

        except requests.exceptions.ConnectionError:
            raise Exception(f"無法連接至 Ollama API ({self.api_url})。\n請確認 Ollama 正在本機端運行。")
        except requests.exceptions.RequestException as e:
            raise Exception(f"請求 Ollama API 時發生錯誤：{e}")
        except Exception as e:
            raise Exception(f"處理 Ollama 回應時發生未知錯誤：{e}")

class PptxService:
    """生成 PowerPoint 投影片"""
    def add_to_presentation(self, pptx_path, genai_text, project_name):
        if not genai_text:
            raise ValueError("報告內容為空，無法生成投影片。")

        prs = Presentation(pptx_path) if os.path.exists(pptx_path) else Presentation()
        reports = re.split(r'(?=--- 報告 \d+[:：])', genai_text)
        generated_count = 0

        for report_text in reports:
            report_text = report_text.strip()
            if not report_text:
                continue
            
            # --- THIS IS THE FIX ---
            # 我們在結尾的 "---" 前後加上 \s* 來允許任意數量的空格
            match = re.match(r'--- 報告 \d+[:：](.*?)\s*---\s*\n(.*)', report_text, re.DOTALL)
            # --- END OF FIX ---

            if not match:
                continue

            ai_title = match.group(1).strip()
            content = match.group(2).strip()
            final_title = f"{project_name} - {ai_title}" if project_name else ai_title

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = final_title
            
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            text_frame.word_wrap = True

            for line in content.split('\n'):
                original_line = line.strip()
                if not original_line: continue

                cleaned_text = re.sub(r'^\s*[-*•\s`]+', '', original_line).strip().strip('*').strip()
                if not cleaned_text: continue

                is_star_heading = any(cleaned_text.lower().startswith(keyword.lower()) for keyword in STAR_KEYWORDS)

                p = text_frame.add_paragraph()
                p.text = cleaned_text
                if is_star_heading:
                    p.level = 0
                    p.font.bold = True
                    p.font.size = Pt(18)
                else:
                    p.level = 1
                    p.font.bold = False
                    p.font.size = Pt(16)
            
            generated_count += 1
        
        if generated_count > 0:
            prs.save(pptx_path)
            return generated_count
        else:
            raise ValueError("未能在內容中找到符合格式的報告。")

