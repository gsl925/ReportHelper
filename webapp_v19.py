import streamlit as st
from PIL import Image
import pytesseract
import extract_msg
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import re
import io

import cv2
import numpy as np

from streamlit_paste_button import paste_image_button

# --- 設定 Tesseract 路徑 ---
if os.name == 'nt':
    try:
        pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        pytesseract.get_tesseract_version()
    except Exception:
        st.warning("Tesseract OCR 未在預設路徑找到，請確認已安裝並設定好環境變數。OCR 功能可能無法使用。")

# --- PROMPTS UPDATED FOR v15.3 ---
SINGLE_ISSUE_PROMPT = """
# Role and Goal
You are a senior technical analyst. Your goal is to analyze the following text, assuming it describes ONE SINGLE core problem, and synthesize all information into ONE extremely concise, direct, and non-repetitive STAR method report.

# Core Principle 1: Brevity is Key
- **Be Direct**: Avoid verbose corporate language. Use clear, simple, and direct technical phrasing.
- **Example**: Instead of "找出...的根本原因，提出有效解決方案...", write "找到...的當機原因。".

# Core Principle 2: Component & Version Identification
- **Identify and Integrate**: You must identify any specific application, program, or component names AND their associated version numbers if provided (e.g., 'v1.2', 'build 22H2', 'rev. A', 'R01A版').
- **Critical Context**: This combination of name and version is crucial. They MUST be explicitly mentioned together in the **Situation (情境)** section.
- **Example**: If the text mentions "OurApp v2.1 fails", the 'Situation' must state "應用程式 OurApp v2.1 發生錯誤...".

# Core Principle 3: Special Data Handling
- **Failure Rate is Priority**: If you identify a failure rate (e.g., 'X/Y failed', 'Z% fail rate', '不良率', '再現率'), you MUST synthesize and include this data point in the 'Situation' section. Highlight it with a bold title like "**Fail Rate:**".

# Core Principle 4: Technical Terminology Preservation
- **Preserve Original Terms**: You MUST preserve original English technical terms, keywords, and proper nouns found in the source text. DO NOT translate them into Chinese. This is critical for technical accuracy.
- **Examples**:
  - If the text says "the test will fail", your report should use "測試會 fail", NOT "測試會失敗".
  - If the text mentions "check the power status", your report should use "檢查 power 狀態", NOT "檢查電源狀態".
- **Scope**: This applies to all technical jargon, acronyms (e.g., BIOS, POST), component names (e.g., CPU, DIMM), status words (e.g., pass, fail, error), and specific commands or values.

# Core Principle 5: Causal and Status Analysis
- **Distinguish Root Causes**: If the text describes multiple distinct root causes for the same high-level problem (e.g., one machine fails due to power, another due to thermal), you MUST describe them as separate findings.
- **Separate Past from Future**: You MUST differentiate between:
  1.  **Completed Diagnostic Actions**: Things that have already been done (e.g., "更換...後", "驗證發現").
  2.  **Planned Next Steps**: Future actions or solutions (e.g., "待RD確認", "請RD驗證", "計畫修改").
- **Link Action to Result**: For each completed action, you MUST state its specific outcome (e.g., "更換 thermal module 後 -> 依舊 Fail", "重跑後 -> PASS"). Do not misattribute results.

# Core Principle 6: Redefining the STAR Categories (Advanced Version)
- **情境 (Situation)**: State the main problem, affected components (with versions), and failure rate.
- **任務 (Task)**: State the core objective (e.g., "診斷問題原因並提出解決方案").
- **行動 (Action)**: List ONLY the **completed diagnostic actions** and their direct results.
  - **Correct Example**: "更換 thermal module -> 依舊 Fail."
  - **Incorrect Example**: "修改 power limit table." (This is a future plan, not a completed action).
- **結果 (Result)**: Summarize the key findings from the actions. Most importantly, list the **planned next steps or solutions**.
  - **Example**: "發現更換散熱模組無效。下一步：POWER 團隊將修改 power limit table 並釋出 test bios 供 RD 驗證。"

# Task
1.  **Synthesize**: Read the entire text to understand the single problem, its distinct causes, and the status of all actions.
2.  **Structure**: Organize the facts into a single, concise STAR report, following ALL principles, especially the separation of past actions and future plans.
3.  **Project Name**: The user has provided the project name separately. DO NOT include it in the report title.

# Output Format and Rules
- The entire output MUST be in **Traditional Chinese (繁體中文)**.
- The report title MUST be a short, direct summary of the problem itself. Example: "WINPE下讀取BIOS資訊當機".
- The output must contain only ONE report block, starting with `--- 報告 1：[問題簡述] ---`.
- For all bullet points under the STAR categories, you MUST indent them with two spaces.

{PROJECT_NAME_HOLDER}

Now, analyze the following text as a single problem and generate one concise report, adhering to all principles.
"""

MULTI_ISSUE_PROMPT = """
# Role and Goal
You are a senior technical analyst. Your goal is to analyze a work discussion, identify ALL distinct core problems, and generate a separate, extremely concise, direct, and non-repetitive STAR method report for EACH core problem.

# Core Principle 1: Brevity is Key
- **Be Direct**: Avoid verbose corporate language. Use clear, simple, and direct technical phrasing.
- **Example**: Instead of "找出...的根本原因，提出有效解決方案...", write "找到...的當機原因。".

# Core Principle 2: Component & Version Identification
- **Identify and Integrate**: You must identify any specific application, program, or component names AND their associated version numbers if provided (e.g., 'v1.2', 'build 22H2', 'rev. A', 'R01A版').
- **Critical Context**: This combination of name and version is crucial. They MUST be explicitly mentioned together in the **Situation (情境)** section.
- **Example**: If the text mentions "OurApp v2.1 fails", the 'Situation' must state "應用程式 OurApp v2.1 發生錯誤...".

# Core Principle 3: Special Data Handling
- **Failure Rate is Priority**: If you identify a failure rate (e.g., 'X/Y failed', 'Z% fail rate', '不良率', '再現率'), you MUST synthesize and include this data point in the 'Situation' section of the corresponding report. Highlight it with a bold title like "**Fail Rate:**".

# Core Principle 4: Technical Terminology Preservation
- **Preserve Original Terms**: You MUST preserve original English technical terms, keywords, and proper nouns found in the source text. DO NOT translate them into Chinese. This is critical for technical accuracy.
- **Examples**:
  - If the text says "the test will fail", your report should use "測試會 fail", NOT "測試會失敗".
  - If the text mentions "check the power status", your report should use "檢查 power 狀態", NOT "檢查電源狀態".
- **Scope**: This applies to all technical jargon, acronyms (e.g., BIOS, POST), component names (e.g., CPU, DIMM), status words (e.g., pass, fail, error), and specific commands or values.

# Core Principle 5: Causal and Status Analysis
- **Distinguish Root Causes**: If the text describes multiple distinct root causes for the same high-level problem (e.g., one machine fails due to power, another due to thermal), you MUST describe them as separate findings.
- **Separate Past from Future**: You MUST differentiate between:
  1.  **Completed Diagnostic Actions**: Things that have already been done (e.g., "更換...後", "驗證發現").
  2.  **Planned Next Steps**: Future actions or solutions (e.g., "待RD確認", "請RD驗證", "計畫修改").
- **Link Action to Result**: For each completed action, you MUST state its specific outcome (e.g., "更換 thermal module 後 -> 依舊 Fail", "重跑後 -> PASS"). Do not misattribute results.

# Core Principle 6: Redefining the STAR Categories (Advanced Version)
- **情境 (Situation)**: State the main problem, affected components (with versions), and failure rate.
- **任務 (Task)**: State the core objective (e.g., "診斷問題原因並提出解決方案").
- **行動 (Action)**: List ONLY the **completed diagnostic actions** and their direct results.
  - **Correct Example**: "更換 thermal module -> 依舊 Fail."
  - **Incorrect Example**: "修改 power limit table." (This is a future plan, not a completed action).
- **結果 (Result)**: Summarize the key findings from the actions. Most importantly, list the **planned next steps or solutions**.
  - **Example**: "發現更換散熱模組無效。下一步：POWER 團隊將修改 power limit table 並釋出 test bios 供 RD 驗證。"

# Task
1.  **Identify & Group**: Identify the core problems and group all related messages.
2.  **Synthesize and Analyze**: For EACH core problem, create a single, concise STAR report, following all principles, especially the separation of past actions and future plans.
3.  **Project Name**: The user has provided the project name separately. DO NOT include it in the report titles.

# Output Format and Rules
- The entire output MUST be in **Traditional Chinese (繁體中文)**.
- For each report, the title MUST be a short, direct summary of the problem itself. Example: "WINPE下讀取BIOS資訊當機".
- Use a separator and a composite title: `--- 報告 1：[問題簡述] ---`.
- For all bullet points under the STAR categories, you MUST indent them with two spaces.

{PROJECT_NAME_HOLDER}

Now, analyze the following text and generate all reports in the specified format, adhering to all principles.
"""

# --- 核心邏輯函式 (不變) ---
def process_image_content(image_input):
    try:
        if isinstance(image_input, Image.Image): image = image_input
        else: image = Image.open(image_input)
        open_cv_image = np.array(image.convert('RGB'))
        open_cv_image = open_cv_image[:, :, ::-1].copy()
        gray = cv2.cvtColor(open_cv_image, cv2.COLOR_BGR2GRAY)
        h, w = gray.shape
        corners = [gray[0:10, 0:10], gray[0:10, w-10:w], gray[h-10:h, 0:10], gray[h-10:h, w-10:w]]
        corner_mean = np.mean([np.mean(c) for c in corners])
        if corner_mean > 128:
            st.info("偵測到淺色模式，使用標準二值化。")
            _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        else:
            st.info("偵測到深色模式，使用反向二值化。")
            _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
        custom_config = r'--oem 3 --psm 6'
        text = pytesseract.image_to_string(thresh, lang='chi_tra+chi_sim+eng', config=custom_config)
        if not text.strip(): st.warning("OCR 引擎未能辨識出任何文字。")
        return text
    except pytesseract.TesseractNotFoundError:
        st.error("Tesseract OCR 引擎未找到或路徑錯誤。")
        return ""
    except Exception as e:
        st.error(f"圖片辨識時發生錯誤: {e}")
        st.exception(e)
        return ""

def process_text_content(text_file):
    try:
        return text_file.getvalue().decode("utf-8")
    except Exception as e:
        st.error(f"讀取文字檔時發生錯誤: {e}")
        return ""

def process_msg_content(msg_file):
    try:
        msg = extract_msg.Message(msg_file)
        return f"寄件人：{msg.sender}\n主旨：{msg.subject}\n\n--- 內文 ---\n{msg.body}"
    except Exception as e:
        st.error(f"解析 Email (.msg) 檔案時發生錯誤: {e}")
        return ""

def generate_powerpoint_in_memory(genai_text, project_name, template_file=None):
    try:
        if template_file: prs = Presentation(template_file)
        else: prs = Presentation()
        try: slide_layout = prs.slide_layouts[1]
        except IndexError:
            st.warning("在範本中找不到標準的『標題及內容』版面配置 (索引 1)，將使用第一個可用的版面。")
            slide_layout = prs.slide_layouts[0]
        reports = re.split(r'(?=--- 報告 \d+：)', genai_text)
        generated_count = 0
        for report_text in reports:
            report_text = report_text.strip()
            if not report_text: continue
            match = re.match(r'--- 報告 \d+：(.*?) ---\n(.*)', report_text, re.DOTALL)
            if not match: continue
            ai_title = match.group(1).strip()
            content = match.group(2).strip()
            final_title = f"{project_name} - {ai_title}" if project_name else ai_title
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title: slide.shapes.title.text = final_title
            if not slide.placeholders or len(slide.placeholders) < 2:
                 st.error(f"錯誤：選擇的投影片版面 '{slide_layout.name}' 沒有足夠的內容佔位符。")
                 return None
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear(); text_frame.word_wrap = True
            for line in content.split('\n'):
                if not line.strip(): continue
                is_indented = line.startswith('  ')
                level = 1 if is_indented else 0
                cleaned_line = re.sub(r'^\s*[-*]\s+', '', line).strip()
                p = text_frame.add_paragraph()
                p.text = cleaned_line; p.level = level
                if level == 0: p.font.bold = True; p.font.size = Pt(18)
                else: p.font.bold = False; p.font.size = Pt(16)
            generated_count += 1
        if generated_count == 0:
            st.warning("未能在貼上的內容中找到符合格式的報告。")
            return None
        ppt_io = io.BytesIO()
        prs.save(ppt_io); ppt_io.seek(0)
        return ppt_io
    except Exception as e:
        st.error(f"生成 PowerPoint 時發生錯誤: {e}")
        st.exception(e)
        return None

def handle_file_upload():
    uploaded_file = st.session_state.file_uploader_key
    if uploaded_file is not None:
        with st.spinner('正在處理檔案...'):
            file_ext = os.path.splitext(uploaded_file.name)[1].lower()
            content = ""
            if file_ext in ['.png', '.jpg', '.jpeg']: content = process_image_content(uploaded_file)
            elif file_ext == '.txt': content = process_text_content(uploaded_file)
            elif file_ext == '.msg': content = process_msg_content(uploaded_file)
            st.session_state.ocr_text = content
            st.session_state.full_prompt = ""
        st.success("檔案處理完成！")

# --- Streamlit UI ---
st.set_page_config(layout="wide", page_title="報告整理小幫手 Web")
st.title("報告整理小幫手 (網頁版)")
st.caption("一個將零散問題快速轉換為標準化 STAR 報告的工具")

# --- 初始化 session_state ---
if 'ocr_text' not in st.session_state:
    st.session_state.ocr_text = ""
if 'full_prompt' not in st.session_state:
    st.session_state.full_prompt = ""
if 'ppt_template_file' not in st.session_state:
    st.session_state.ppt_template_file = None

col1, col2 = st.columns(2)

# --- 左側欄：資料輸入與 Prompt 生成 ---
with col1:
    st.header("步驟 A: 輸入原始資料")
    project_name = st.text_input("專案名稱 (選填)", placeholder="例如：BIOS/記憶體專案")
    
    st.file_uploader("選項 1: 上傳檔案 (清空現有內容)", type=['png', 'jpg', 'jpeg', 'txt', 'msg'], key="file_uploader_key", on_change=handle_file_upload)
    
    st.write("選項 2: 從剪貼簿連續貼上圖片")
    
    paste_info = paste_image_button(
        label="📋 附加截圖 (Ctrl+V)", 
        key="paste", 
        background_color="#FF4B4B", 
        hover_background_color="#FF6B6B"
    )

    if paste_info and paste_info.image_data is not None:
        with st.spinner('正在進行 OCR 辨識...'):
            new_content = process_image_content(paste_info.image_data)
            
            if st.session_state.ocr_text.strip():
                separator = f"\n\n{'='*20} 來自剪貼簿的新增圖片 {'='*20}\n\n"
                st.session_state.ocr_text += separator + new_content
            else:
                st.session_state.ocr_text = new_content
            
            st.session_state.full_prompt = ""
        
        st.success("圖片辨識完成，並已附加至結果中！")

    st.text_area("辨識/解析結果", height=250, key="ocr_text")
    
    st.header("步驟 B: 準備分析指令 (Prompt)")
    prompt_col1, prompt_col2 = st.columns(2)
    with prompt_col1:
        if st.button("準備「單一問題」Prompt", use_container_width=True):
            if st.session_state.ocr_text:
                project_info = f"The user has specified the project name is: '{project_name}'." if project_name else "The user did not specify a project name."
                final_prompt = SINGLE_ISSUE_PROMPT.replace("{PROJECT_NAME_HOLDER}", project_info)
                st.session_state.full_prompt = f"{final_prompt}\n\n{st.session_state.ocr_text}"
            else: st.warning("請先上傳檔案或貼上圖片以取得原始文字。")
    with prompt_col2:
        if st.button("準備「多個問題」Prompt", use_container_width=True):
            if st.session_state.ocr_text:
                project_info = f"The user has specified the project name is: '{project_name}'." if project_name else "The user did not specify a project name."
                final_prompt = MULTI_ISSUE_PROMPT.replace("{PROJECT_NAME_HOLDER}", project_info)
                st.session_state.full_prompt = f"{final_prompt}\n\n{st.session_state.ocr_text}"
            else: st.warning("請先上傳檔案或貼上圖片以取得原始文字。")
            
    if st.session_state.full_prompt:
        st.text_area("複製以下完整指令到 GenAI 工具中", value=st.session_state.full_prompt, height=200)

# --- 右側欄：結果貼上與報告生成 (不變) ---
with col2:
    st.header("步驟 C: 貼上分析結果")
    genai_output = st.text_area(
        "在此貼上 GenAI 產生的 STAR 報告",
        height=450,
        key="genai_output_area",
        placeholder="--- 報告 1：[問題簡述] ---\n- 情境 (Situation)\n  - ...\n- 任務 (Task)\n  - ...\n..."
    )
    
    st.header("步驟 D: 產生並下載報告")

    st.subheader("選項：使用現有簡報範本")
    st.session_state.ppt_template_file = st.file_uploader(
        "上傳您的 .pptx 範本 (選填)",
        type=['pptx'],
        key="ppt_template_uploader"
    )
    
    if st.session_state.ppt_template_file:
        st.info("已上傳範本。新報告將會被新增到此檔案的末尾。")
    
    if st.button("產生 PowerPoint 報告", type="primary", use_container_width=True):
        if genai_output:
            with st.spinner("正在生成 PowerPoint 檔案..."):
                ppt_file_in_memory = generate_powerpoint_in_memory(
                    genai_output, 
                    project_name, 
                    st.session_state.ppt_template_file
                )
            
            if ppt_file_in_memory:
                st.success("PowerPoint 報告已生成！")
                
                is_template_used = st.session_state.ppt_template_file is not None
                download_filename = f"{project_name or 'Report'}_Updated.pptx" if is_template_used else f"{project_name or 'Report'}.pptx"

                st.download_button(
                    label="📥 下載報告 (.pptx)",
                    data=ppt_file_in_memory,
                    file_name=download_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
        else:
            st.warning("請先在上方貼上 GenAI 產生的報告內容。")

