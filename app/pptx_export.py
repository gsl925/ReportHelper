# app/pptx_export.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def create_star_slide(prs: Presentation, title: str, star: dict):
    slide_layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    def add_section(y, header, items):
        box = slide.shapes.add_textbox(left, Inches(y), Inches(9), Inches(1.8))
        tf = box.text_frame
        h = tf.paragraphs[0]
        h.text = header
        h.font.size = Pt(16)
        h.font.bold = True
        for it in items:
            p = tf.add_paragraph()
            p.text = '• ' + it
            p.level = 1
            p.font.size = Pt(14)

    add_section(1.3, 'Situation', star.get('situation', []))
    add_section(3.0, 'Task', star.get('task', []))
    add_section(4.3, 'Action', star.get('action', []))
    add_section(6.0, 'Result', star.get('result', []))


def export_to_pptx(filename: str, title: str, star: dict, author: str = 'Report Helper'):
    prs = Presentation()
    create_star_slide(prs, title, star)
    prs.core_properties.author = author
    prs.save(filename)
    return filename