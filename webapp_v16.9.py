import streamlit as st
from PIL import Image
import pytesseract
import extract_msg
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import re
import io

# NEW: 引入 OpenCV 和 NumPy
import cv2
import numpy as np

from streamlit_paste_button import paste_image_button

# --- 設定 Tesseract 路徑 (與之前相同) ---
if os.name == 'nt':
    try:
        pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        pytesseract.get_tesseract_version()
    except Exception:
        st.warning("Tesseract OCR 未在預設路徑找到，請確認已安裝並設定好環境變數。OCR 功能可能無法使用。")

# --- FINAL CHANGE: 這次包含完整的 Prompt 內容 ---
SINGLE_ISSUE_PROMPT = """
# Role and Goal
You are a senior technical analyst. Your goal is to analyze the following text, assuming it describes ONE SINGLE core problem, and synthesize all information into ONE extremely concise, direct, and non-repetitive STAR method report.
# Core Principle 1: Brevity is Key
- **Be Direct**: Avoid verbose corporate language. Use clear, simple, and direct technical phrasing.
- **Example**: Instead of "找出...的根本原因，提出有效解決方案...", write "找到...的當機原因。".
# Core Principle 2: Component & Version Identification
- **Identify and Integrate**: You must identify any specific application, program, or component names AND their associated version numbers if provided (e.g., 'v1.2', 'build 22H2', 'rev. A', 'R01A版').
- **Critical Context**: This combination of name and version is crucial. They MUST be explicitly mentioned together in the **Situation (情境)** section.
- **Example**: If the text mentions "OurApp v2.1 fails", the 'Situation' must state "應用程式 OurApp v2.1 發生錯誤...".
# Core Principle 3: Special Data Handling
- **Failure Rate is Priority**: If you identify a failure rate (e.g., 'X/Y failed', 'Z% fail rate', '不良率', '再現率'), you MUST synthesize and include this data point in the 'Situation' section. Highlight it with a bold title like "**問題再現率 (Fail Rate):**".
# Core Principle 4: Redefining the STAR Categories (Concise Version)
- **情境 (Situation)**: Directly state the problem, including component names and their versions. Must include failure rate if available.
- **任務 (Task)**: State the core objective in a few words (e.g., "診斷問題原因").
- **行動 (Action)**: List key diagnostic steps. No extra descriptions.
- **結果 (Result)**: Summarize key findings and next steps. Do not repeat facts from other sections.
# Task
1.  **Synthesize**: Read the entire text to understand the single problem.
2.  **Structure**: Organize the facts into a single, concise STAR report, following all principles.
3.  **Project Name**: The user has provided the project name separately. DO NOT include it in the report title.
# Output Format and Rules
- The entire output MUST be in **Traditional Chinese (繁體中文)**.
- The report title MUST be a short, direct summary of the problem itself. Example: "WINPE下讀取BIOS資訊當機".
- The output must contain only ONE report block, starting with `--- 報告 1：[問題簡述] ---`.
- For all bullet points under the STAR categories, you MUST indent them with two spaces.
{PROJECT_NAME_HOLDER}
Now, analyze the following text as a single problem and generate one concise report, adhering to all principles.
"""

MULTI_ISSUE_PROMPT = """
# Role and Goal
You are a senior technical analyst. Your goal is to analyze a work discussion, identify ALL distinct core problems, and generate a separate, extremely concise, direct, and non-repetitive STAR method report for EACH core problem.
# Core Principle 1: Brevity is Key
- **Be Direct**: Avoid verbose corporate language. Use clear, simple, and direct technical phrasing.
- **Example**: Instead of "找出...的根本原因，提出有效解決方案...", write "找到...的當機原因。".
# Core Principle 2: Component & Version Identification
- **Identify and Integrate**: You must identify any specific application, program, or component names AND their associated version numbers if provided (e.g., 'v1.2', 'build 22H2', 'rev. A', 'R01A版').
- **Critical Context**: This combination of name and version is crucial. They MUST be explicitly mentioned together in the **Situation (情境)** section.
- **Example**: If the text mentions "OurApp v2.1 fails", the 'Situation' must state "應用程式 OurApp v2.1 發生錯誤...".
# Core Principle 3: Special Data Handling
- **Failure Rate is Priority**: If you identify a failure rate (e.g., 'X/Y failed', 'Z% fail rate', '不良率', '再現率'), you MUST synthesize and include this data point in the 'Situation' section of the corresponding report. Highlight it with a bold title like "**問題再現率 (Fail Rate):**".
# Core Principle 4: Redefining the STAR Categories (Concise Version)
- **情境 (Situation)**: Directly state the problem, including component names and their versions. Must include failure rate if available.
- **任務 (Task)**: State the core objective in a few words (e.g., "診斷問題原因").
- **行動 (Action)**: List key diagnostic steps. No extra descriptions.
- **結果 (Result)**: Summarize key findings and next steps. Do not repeat facts from other sections.
# Task
1.  **Identify & Group**: Identify the core problems and group all related messages.
2.  **Synthesize and Analyze**: For EACH core problem, create a single, concise STAR report, following all principles.
3.  **Project Name**: The user has provided the project name separately. DO NOT include it in the report titles.
# Output Format and Rules
- The entire output MUST be in **Traditional Chinese (繁體中文)**.
- For each report, the title MUST be a short, direct summary of the problem itself. Example: "WINPE下讀取BIOS資訊當機".
- Use a separator and a composite title: `--- 報告 1：[問題簡述] ---`.
- For all bullet points under the STAR categories, you MUST indent them with two spaces.
{PROJECT_NAME_HOLDER}
Now, analyze the following text and generate all reports in the specified format, adhering to all principles.
"""

# --- 核心邏輯函式 ---

# FINAL, MAJOR CHANGE: 引入 OpenCV 進行影像前處理，大幅提升 OCR 準確率
def process_image_content(image_input):
    """
    使用 Tesseract 進行 OCR，並在之前進行影像前處理。
    """
    try:
        if isinstance(image_input, Image.Image):
            image = image_input
        else:
            image = Image.open(image_input)
        
        # 1. 將 PIL 圖片物件轉換為 OpenCV 格式 (NumPy array)
        #    注意：PIL 的 RGB 和 OpenCV 的 BGR 顏色通道順序相反
        open_cv_image = np.array(image.convert('RGB'))
        open_cv_image = open_cv_image[:, :, ::-1].copy()

        # 2. 灰階化 (Grayscaling)
        gray = cv2.cvtColor(open_cv_image, cv2.COLOR_BGR2GRAY)

        # 3. 二值化 (Binarization) - 這是最關鍵的步驟
        #    我們使用 Otsu's 方法自動尋找最佳閾值
        _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

        # 4. (可選) 放大圖片，有時能提升對小字體的辨識率
        # h, w = thresh.shape
        # if h < 100 or w < 100: # 如果圖片太小
        #     thresh = cv2.resize(thresh, (w*2, h*2), interpolation=cv2.INTER_CUBIC)

        # 5. 將處理後的圖片傳給 Tesseract
        #    我們也加入 --psm 6 參數，告知 Tesseract 這可能是一個統一的文字區塊
        custom_config = r'--oem 3 --psm 6'
        text = pytesseract.image_to_string(thresh, lang='chi_tra+chi_sim+eng', config=custom_config)
        
        return text
    except pytesseract.TesseractNotFoundError:
        st.error("Tesseract OCR 引擎未找到或路徑錯誤。請檢查您的伺服器環境設定。")
        return ""
    except Exception as e:
        st.error(f"圖片辨識時發生錯誤: {e}")
        return ""

# --- 其他核心函式 (與之前相同) ---
def process_text_content(text_file):
    # ... (程式碼不變)
    try:
        content = text_file.getvalue().decode("utf-8")
        return content
    except Exception as e:
        st.error(f"讀取文字檔時發生錯誤: {e}")
        return ""

def process_msg_content(msg_file):
    # ... (程式碼不變)
    try:
        msg = extract_msg.Message(msg_file)
        formatted_content = (
            f"寄件人：{msg.sender}\n"
            f"主旨：{msg.subject}\n\n"
            f"--- 內文 ---\n"
            f"{msg.body}"
        )
        return formatted_content
    except Exception as e:
        st.error(f"解析 Email (.msg) 檔案時發生錯誤: {e}")
        return ""

def generate_powerpoint_in_memory(genai_text, project_name):
    # ... (程式碼不變)
    try:
        prs = Presentation()
        reports = re.split(r'(?=--- 報告 \d+：)', genai_text)
        generated_count = 0
        for report_text in reports:
            report_text = report_text.strip()
            if not report_text: continue
            match = re.match(r'--- 報告 \d+：(.*?) ---\n(.*)', report_text, re.DOTALL)
            if not match: continue
            ai_title = match.group(1).strip()
            content = match.group(2).strip()
            final_title = f"{project_name} - {ai_title}" if project_name else ai_title
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = final_title
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            for line in content.split('\n'):
                if not line.strip(): continue
                is_indented = line.startswith('  ')
                level = 1 if is_indented else 0
                cleaned_line = re.sub(r'^\s*[-*]\s+', '', line).strip()
                p = text_frame.add_paragraph()
                p.text = cleaned_line
                p.level = level
                if level == 0:
                    p.font.bold = True
                    p.font.size = Pt(18)
                else:
                    p.font.bold = False
                    p.font.size = Pt(16)
            generated_count += 1
        if generated_count == 0:
            st.warning("未能在貼上的內容中找到符合格式的報告，請檢查內容。")
            return None
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        return ppt_io
    except Exception as e:
        st.error(f"生成 PowerPoint 時發生錯誤: {e}")
        return None

def handle_file_upload():
    # ... (程式碼不變)
    uploaded_file = st.session_state.file_uploader_key
    if uploaded_file is not None:
        with st.spinner('正在處理檔案...'):
            file_ext = os.path.splitext(uploaded_file.name)[1].lower()
            content = ""
            if file_ext in ['.png', '.jpg', '.jpeg']:
                content = process_image_content(uploaded_file)
            elif file_ext == '.txt':
                content = process_text_content(uploaded_file)
            elif file_ext == '.msg':
                content = process_msg_content(uploaded_file)
            st.session_state.ocr_text = content
            st.session_state.full_prompt = ""
        st.success("檔案處理完成！")

# --- Streamlit UI 介面 (與之前相同) ---
st.set_page_config(layout="wide", page_title="報告整理小幫手 Web")
st.title("報告整理小幫手 (網頁版)")
st.caption("一個將零散問題快速轉換為標準化 STAR 報告的工具")

if 'ocr_text' not in st.session_state:
    st.session_state.ocr_text = ""
if 'full_prompt' not in st.session_state:
    st.session_state.full_prompt = ""

col1, col2 = st.columns(2)

with col1:
    st.header("步驟 A: 輸入原始資料")
    project_name = st.text_input("專案名稱 (選填)", placeholder="例如：BIOS/記憶體專案")
    st.file_uploader(
        "選項 1: 拖曳或點擊上傳檔案",
        type=['png', 'jpg', 'jpeg', 'txt', 'msg'],
        key="file_uploader_key",
        on_change=handle_file_upload
    )
    st.write("選項 2: 從剪貼簿貼上圖片")
    paste_info = paste_image_button(
        label="📋 貼上截圖 (Ctrl+V)",
        key="paste",
        background_color="#FF4B4B",
        hover_background_color="#FF6B6B"
    )
    if paste_info and paste_info.image_data is not None:
        st.write("已成功貼上圖片！")
        with st.spinner('正在進行 OCR 辨識...'):
            content = process_image_content(paste_info.image_data)
            st.session_state.ocr_text = content
            st.session_state.full_prompt = ""
        st.success("圖片辨識完成！")
    st.text_area("辨識/解析結果", value=st.session_state.ocr_text, height=250, key="ocr_text_display")
    st.header("步驟 B: 準備分析指令 (Prompt)")
    prompt_col1, prompt_col2 = st.columns(2)
    with prompt_col1:
        if st.button("準備「單一問題」Prompt", use_container_width=True):
            if st.session_state.ocr_text:
                project_info = f"The user has specified the project name is: '{project_name}'." if project_name else "The user did not specify a project name."
                final_prompt = SINGLE_ISSUE_PROMPT.replace("{PROJECT_NAME_HOLDER}", project_info)
                st.session_state.full_prompt = f"{final_prompt}\n\n{st.session_state.ocr_text}"
            else:
                st.warning("請先上傳檔案以取得原始文字。")
    with prompt_col2:
        if st.button("準備「多個問題」Prompt", use_container_width=True):
            if st.session_state.ocr_text:
                project_info = f"The user has specified the project name is: '{project_name}'." if project_name else "The user did not specify a project name."
                final_prompt = MULTI_ISSUE_PROMPT.replace("{PROJECT_NAME_HOLDER}", project_info)
                st.session_state.full_prompt = f"{final_prompt}\n\n{st.session_state.ocr_text}"
            else:
                st.warning("請先上傳檔案以取得原始文字。")
    if st.session_state.full_prompt:
        st.text_area("複製以下完整指令到 GenAI 工具中", value=st.session_state.full_prompt, height=200)

with col2:
    st.header("步驟 C: 貼上分析結果")
    genai_output = st.text_area(
        "在此貼上 GenAI 產生的 STAR 報告",
        height=450,
        key="genai_output_area",
        placeholder="--- 報告 1：[問題簡述] ---\n- 情境 (Situation)\n  - ...\n- 任務 (Task)\n  - ...\n..."
    )
    st.header("步驟 D: 產生並下載報告")
    if st.button("產生 PowerPoint 報告", type="primary", use_container_width=True):
        if genai_output:
            with st.spinner("正在生成 PowerPoint 檔案..."):
                ppt_file_in_memory = generate_powerpoint_in_memory(genai_output, project_name)
            if ppt_file_in_memory:
                st.success("PowerPoint 報告已生成！")
                st.download_button(
                    label="📥 下載報告 (.pptx)",
                    data=ppt_file_in_memory,
                    file_name=f"{project_name or 'Report'}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
        else:
            st.warning("請先在上方貼上 GenAI 產生的報告內容。")
